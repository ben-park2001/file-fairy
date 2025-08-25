"""
File parsing utilities for extracting text content from various file types.
"""

import os
import re
import struct
import unicodedata
import zlib
from pathlib import Path
import fitz

from ..models.schema import ExtractTextRequest, ExtractTextResponse


def extract_text(request: ExtractTextRequest) -> ExtractTextResponse:
    """
    Extract text from a file using the schema-based approach.

    Args:
        request (ExtractTextRequest): The extraction request

    Returns:
        ExtractTextResponse: The structured extraction response
    """
    try:
        result = extract_text_structured(request.file_path, request.clean)
        return ExtractTextResponse(**result)
    except FileNotFoundError:
        return ExtractTextResponse(
            file_name=Path(request.file_path).name,
            file_type=Path(request.file_path).suffix.lower(),
            content="",
            error="File not found",
        )
    except PermissionError:
        return ExtractTextResponse(
            file_name=Path(request.file_path).name,
            file_type=Path(request.file_path).suffix.lower(),
            content="",
            error="Permission denied",
        )
    except ValueError as e:
        return ExtractTextResponse(
            file_name=Path(request.file_path).name,
            file_type=Path(request.file_path).suffix.lower(),
            content="",
            error=str(e),
        )
    except Exception as e:
        return ExtractTextResponse(
            file_name=Path(request.file_path).name,
            file_type=Path(request.file_path).suffix.lower(),
            content="",
            error=f"Extraction failed: {str(e)}",
        )


def get_supported_file_types() -> list[str]:
    """
    Get a list of file extensions that are supported for text extraction.

    Returns:
        list[str]: List of supported file extensions
    """
    return [
        ".txt",
        ".md",
        ".log",
        ".csv",  # Text files
        ".pdf",  # PDF documents
        ".docx",
        ".doc",  # Word documents
        ".xlsx",
        ".xls",  # Excel spreadsheets
        ".pptx",
        ".ppt",  # PowerPoint presentations
        ".hwp",
        ".hwpx",  # Hanword documents
    ]


def is_file_supported(file_path: str) -> bool:
    """
    Check if a file type is supported for text extraction.

    Args:
        file_path (str): Path to the file

    Returns:
        bool: True if the file type is supported
    """
    file_extension = Path(file_path).suffix.lower()
    return file_extension in get_supported_file_types()


def clean_content(text: str) -> str:
    """
    Clean extracted text content by removing unwanted characters and normalizing whitespace.

    Args:
        text (str): Raw extracted text

    Returns:
        str: Cleaned text content
    """
    if not text or not isinstance(text, str):
        return ""

    # Remove control characters except newlines and tabs
    text = "".join(
        char for char in text if unicodedata.category(char)[0] != "C" or char in "\n\t"
    )

    # Remove excessive whitespace while preserving line structure
    text = re.sub(r"[ \t]+", " ", text)  # Multiple spaces/tabs to single space
    text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)  # Multiple newlines to double newline

    # Remove leading/trailing whitespace from each line
    lines = [line.strip() for line in text.split("\n")]
    text = "\n".join(lines)

    # Remove empty lines at start and end
    text = text.strip()

    return text


def extract_text_structured(file_path: str, clean: bool = True) -> dict:
    """
    Extract text from a file and return structured data.

    Args:
        file_path (str): Absolute path to the file
        clean (bool): Whether to clean the extracted content

    Returns:
        dict: Structured extraction results
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    if not os.access(file_path, os.R_OK):
        raise PermissionError(f"Cannot read file: {file_path}")

    path = Path(file_path)
    file_extension = path.suffix.lower()

    # Determine extraction method and extract content
    if file_extension in [".txt", ".md", ".log", ".csv"]:
        content = _extract_text_file(file_path)
    elif file_extension == ".pdf":
        content = _extract_pdf_file(file_path)
    elif file_extension in [".docx", ".doc"]:
        content = _extract_word_file(file_path)
    elif file_extension in [".xlsx", ".xls"]:
        content = _extract_excel_file(file_path)
    elif file_extension in [".pptx", ".ppt"]:
        content = _extract_powerpoint_file(file_path)
    elif file_extension in [".hwp", ".hwpx"]:
        content = _extract_hwp_file(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

    # Clean content if requested
    if clean:
        content = clean_content(content)

    return {
        "success": True,
        "file_name": path.name,
        "file_type": file_extension,
        "content": content,
    }


def _extract_text_file(file_path: str) -> str:
    """Extract text from plain text files."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read()
        except Exception:
            return ""


def _extract_pdf_file(file_path: str) -> str:
    """Extract text from PDF files."""
    try:
        doc = fitz.open(file_path)
        text = ""
        for page_num in range(len(doc)):
            page = doc[page_num]
            text += page.get_text()
        doc.close()
        return text
    except Exception:
        return ""


def _extract_word_file(file_path: str) -> str:
    """Extract text from Word documents."""
    try:
        import docx2txt

        return docx2txt.process(file_path) or ""
    except Exception:
        try:
            from docx import Document

            doc = Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception:
            return ""


def _extract_excel_file(file_path: str) -> str:
    """Extract text from Excel files."""
    try:
        file_extension = Path(file_path).suffix.lower()
        text_content = []

        if file_extension == ".xlsx":
            import openpyxl

            workbook = openpyxl.load_workbook(file_path, data=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows(values=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text_content.append(" ".join(row_text))
            workbook.close()
        elif file_extension == ".xls":
            import xlrd

            workbook = xlrd.open_workbook(file_path)
            for sheet_name in workbook.sheet_names():
                sheet = workbook.sheet_by_name(sheet_name)
                for row_idx in range(sheet.nrows):
                    row_text = [
                        str(sheet.cell(row_idx, col_idx).value)
                        for col_idx in range(sheet.ncols)
                        if sheet.cell(row_idx, col_idx).value
                    ]
                    if row_text:
                        text_content.append(" ".join(row_text))

        return "\n".join(text_content)
    except Exception:
        return ""


def _extract_powerpoint_file(file_path: str) -> str:
    """Extract text from PowerPoint files."""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_content = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                text_content.append(run.text)
                elif hasattr(shape, "text"):
                    shape_text = getattr(shape, "text", "")
                    if shape_text.strip():
                        text_content.append(shape_text)

        return "\n".join(text_content)
    except Exception:
        return ""


def _extract_hwp_file(file_path: str) -> str:
    """Main HWP extraction logic."""
    try:
        import olefile
    except ImportError as e:
        raise ImportError(f"Required libraries not available for HWP extraction: {e}")

    # HWP file constants
    FILE_HEADER_SECTION = "FileHeader"
    HWP_SUMMARY_SECTION = "\x05HwpSummaryInformation"
    SECTION_NAME_LENGTH = len("Section")
    BODYTEXT_SECTION = "BodyText"
    HWP_TEXT_TAGS = [67]

    ole_file = olefile.OleFileIO(file_path)
    directories = ole_file.listdir()

    # Validate HWP file
    has_header = [FILE_HEADER_SECTION] in directories
    has_summary = [HWP_SUMMARY_SECTION] in directories
    if not (has_header and has_summary):
        ole_file.close()
        raise ValueError("Not a valid HWP file")

    # Check if compressed
    header_stream = ole_file.openstream("FileHeader")
    header_data = header_stream.read()
    is_compressed = (header_data[36] & 1) == 1

    # Get body sections
    section_numbers = []
    for directory in directories:
        if directory[0] == BODYTEXT_SECTION:
            section_num = int(directory[1][SECTION_NAME_LENGTH:])
            section_numbers.append(section_num)

    sections = [f"BodyText/Section{num}" for num in sorted(section_numbers)]

    # Extract text from all sections
    extracted_text = ""
    for section in sections:
        extracted_text += _extract_hwp_section_text(
            ole_file, section, is_compressed, HWP_TEXT_TAGS
        )
        extracted_text += "\n"

    ole_file.close()
    return extracted_text.strip()


def _extract_hwp_section_text(
    ole_file, section_name: str, is_compressed: bool, hwp_text_tags: list
) -> str:
    """Extract text from a specific HWP section."""
    section_stream = ole_file.openstream(section_name)
    raw_data = section_stream.read()

    if is_compressed:
        try:
            unpacked_data = zlib.decompress(raw_data, -15)
        except zlib.error:
            return ""
    else:
        unpacked_data = raw_data

    # Parse section data to extract text content
    size = len(unpacked_data)
    position = 0
    text_content = ""

    while position < size:
        try:
            header = struct.unpack_from("<I", unpacked_data, position)[0]
            record_type = header & 0x3FF
            record_length = (header >> 20) & 0xFFF

            if record_type in hwp_text_tags:
                record_data = unpacked_data[position + 4 : position + 4 + record_length]
                decoded_text = _decode_hwp_record_data(record_data)
                if decoded_text:
                    text_content += decoded_text + "\n"

            position += 4 + record_length

        except (struct.error, IndexError):
            # Skip problematic data and continue
            position += 1
            continue

    return text_content


def _decode_hwp_record_data(record_data: bytes) -> str:
    """Decode HWP record data to text."""
    try:
        decoded_text = record_data.decode("utf-16le")

        # Clean extracted text by removing unwanted characters
        # Remove Chinese characters
        text = re.sub(r"[\u4e00-\u9fff]+", "", decoded_text)

        # Remove control characters
        text = "".join(char for char in text if unicodedata.category(char)[0] != "C")

        # Remove excessive whitespace
        text = re.sub(r"\s+", " ", text).strip()

        return text
    except UnicodeDecodeError:
        return ""
