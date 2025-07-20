"""Extractor for Microsoft Office documents."""

from pathlib import Path
from .base import BaseExtractor


class DOCXExtractor(BaseExtractor):
    """Extractor for DOCX files."""
    
    def extract(self, file_path: Path) -> str:
        """Extract text from DOCX files."""
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [paragraph.text for paragraph in doc.paragraphs]
            return '\n'.join(paragraphs)
        except ImportError:
            return self._handle_import_error("DOCX", "pip install python-docx")
        except Exception as e:
            return self._handle_extraction_error(file_path, e)


class PPTXExtractor(BaseExtractor):
    """Extractor for PPTX files."""
    
    def extract(self, file_path: Path) -> str:
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation
            presentation = Presentation(file_path)
            text_content = []
            
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            
            return '\n'.join(text_content)
        except ImportError:
            return self._handle_import_error("PPTX", "pip install python-pptx")
        except Exception as e:
            return self._handle_extraction_error(file_path, e)
