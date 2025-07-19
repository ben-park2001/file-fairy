# aifo.py
# AI File Organizer CLI - Updated with Gemma ONNX Support
# 
# 필요한 패키지 설치:
# pip install pdfplumber python-docx python-pptx openpyxl olefile

import argparse
import datetime
import os
import re
import shutil
import sys
from pathlib import Path

# --- AI 모듈 Import ---
# ai_module.py가 같은 폴더에 있어야 합니다.
try:
    from ai_module import AIKeywordExtractor
except ImportError:
    print("오류: 'ai_module.py' 파일을 찾을 수 없습니다. 스크립트와 같은 폴더에 있는지 확인해주세요.")
    sys.exit(1)

# --- 설정 (Configuration) ---

# 모델 경로 (사용자가 수정해야 하는 부분)
DEFAULT_MODEL_PATH = "gemma-3n-E2B-it-ONNX"

# 기본 프롬프트 템플릿
DEFAULT_PROMPT_TEMPLATE = """You are an expert file naming specialist. Your sole task is to generate new filenames based on file content and original filename.

TASK: Generate a single, optimized filename that accurately represents the file's content.

GUIDELINES:
- Length: 3-8 words maximum
- Format: Use underscores instead of spaces
- Include file extension from original filename
- Be descriptive but concise
- Use common, clear terminology
- Avoid special characters except underscores and periods

INPUT FORMAT:
Original filename: [filename]
File content: [content]

OUTPUT FORMAT:
Provide ONLY the new filename. No explanations, no additional text.

EXAMPLES:
Original filename: document1.pdf
File content: This is a quarterly sales report for Q3 2024 showing revenue growth of 15%
Output: Q3_2024_sales_report.pdf

Original filename: IMG_20240315.jpg
File content: Photo of a golden retriever playing in a park during sunset
Output: golden_retriever_sunset_park.jpg

Original filename: notes.txt
File content: Meeting notes from project kickoff discussing timeline, budget allocation, and team responsibilities
Output: project_kickoff_meeting_notes.txt

Now process the following file:
---
{file_content}
---
"""

# 파일 분류 규칙
FILE_CATEGORIES = {
    '문서': ['.pdf', '.docx', '.pptx', '.hwp', '.txt', '.md'],
    '스프레드시트': ['.xlsx', '.xls', '.csv'],
    '이미지': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff'],
    '동영상': ['.mp4', '.avi', '.mkv', '.mov', '.wmv'],
    '음악': ['.mp3', '.wav', '.flac', '.m4a'],
    '압축파일': ['.zip', '.7z', '.rar', '.tar', '.gz'],
    '코드': ['.py', '.js', '.html', '.css', '.java', '.cpp', '.c'],
    '기타': 'others'
}

# AI 처리 대상 파일 확장자
AI_SUPPORTED_EXTENSIONS = ['.txt', '.md', '.pdf', '.docx', '.pptx', '.hwp', '.xlsx', '.xls']
# 이미지 AI 처리 대상 파일 확장자
AI_SUPPORTED_IMAGE_EXTENSIONS = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']
# 오디오 AI 처리 대상 파일 확장자
AI_SUPPORTED_AUDIO_EXTENSIONS = ['.mp3', '.wav', '.flac', '.m4a', '.ogg', '.aac']

# 제외할 파일/폴더 목록
EXCLUDE_PATTERNS = [
    '.git', '.DS_Store', 'node_modules', 'venv', '__pycache__',
    'Thumbs.db', '.vscode', '.idea', 'aifo.log'
]

# 로그 파일명
LOG_FILE_NAME = "aifo.log"

# AI Extractor 인스턴스를 전역으로 관리 (Lazy Loading)
AI_EXTRACTOR: AIKeywordExtractor | None = None

# --- 문서 읽기 함수 ---

def extract_text_from_file(file_path: Path) -> str:
    """파일 형식에 따라 적절한 방법으로 텍스트를 추출합니다."""
    ext = file_path.suffix.lower()
    
    try:
        if ext == '.pdf':
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    text = ""
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    return text
            except ImportError:
                print("경고: pdfplumber가 설치되지 않았습니다. pip install pdfplumber로 설치해주세요.")
                return ""
                
        elif ext == '.docx':
            try:
                from docx import Document
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            except ImportError:
                print("경고: python-docx가 설치되지 않았습니다. pip install python-docx로 설치해주세요.")
                return ""
            
        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text += shape.text + "\n"
                        # 표 내용도 추출
                        if shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if cell.text.strip():
                                        text += cell.text + " "
                                text += "\n"
                return text
            except ImportError:
                print("경고: python-pptx가 설치되지 않았습니다. pip install python-pptx로 설치해주세요.")
                return ""
            except Exception as e:
                print(f"경고: PPTX 파일 '{file_path.name}' 읽기에 실패했습니다: {e}")
                return ""
                
        elif ext in ['.xlsx', '.xls']:
            try:
                import openpyxl
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                text = ""
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text += f"[시트: {sheet_name}]\n"
                    for row in sheet.iter_rows(values_only=True):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                row_text.append(str(cell))
                        if row_text and any(cell.strip() for cell in row_text if isinstance(cell, str)):
                            text += " | ".join(row_text) + "\n"
                    text += "\n"
                return text
            except ImportError:
                print("경고: openpyxl이 설치되지 않았습니다. pip install openpyxl로 설치해주세요.")
                return ""
            except Exception as e:
                print(f"경고: Excel 파일 '{file_path.name}' 읽기에 실패했습니다: {e}")
                return ""
                
        elif ext == '.hwp':
            try:
                # 먼저 olefile을 사용해서 HWP 파일 읽기 시도
                import olefile
                if olefile.isOleFile(file_path):
                    ole = olefile.OleFileIO(file_path)
                    # HWP 파일의 텍스트 스트림 찾기
                    if ole._olestream_size is not None:
                        # 간단한 텍스트 추출 (완전하지 않을 수 있음)
                        with open(file_path, 'rb') as f:
                            data = f.read()
                            # 한글 텍스트 패턴 찾기 (매우 기본적인 방법)
                            text_parts = []
                            try:
                                # UTF-16으로 디코딩 시도
                                decoded = data.decode('utf-16le', errors='ignore')
                                # 의미있는 텍스트만 추출
                                import re
                                korean_text = re.findall(r'[가-힣\s]+', decoded)
                                text_parts.extend([t.strip() for t in korean_text if len(t.strip()) > 2])
                            except:
                                pass
                            
                            try:
                                # UTF-8로도 시도
                                decoded = data.decode('utf-8', errors='ignore')
                                korean_text = re.findall(r'[가-힣\s]+', decoded)
                                text_parts.extend([t.strip() for t in korean_text if len(t.strip()) > 2])
                            except:
                                pass
                            
                        ole.close()
                        return "\n".join(text_parts[:100])  # 상위 100개 텍스트 부분만
                    ole.close()
            except ImportError:
                print("경고: olefile이 설치되지 않았습니다. pip install olefile로 설치해주세요.")
                return ""
            except Exception as e:
                print(f"경고: HWP 파일 '{file_path.name}' 읽기에 실패했습니다: {e}")
                return ""
            return ""
                
        elif ext in ['.txt', '.md']:
            return file_path.read_text(encoding='utf-8', errors='ignore')
            
        else:
            # 지원하지 않는 형식은 텍스트로 읽기 시도
            try:
                return file_path.read_text(encoding='utf-8', errors='ignore')
            except:
                return ""
            
    except Exception as e:
        print(f"파일 읽기 오류 ({file_path.name}): {e}")
        return ""

# --- 유틸리티 함수 ---

def get_file_category(file_ext: str) -> str:
    """파일 확장자에 따른 카테고리를 반환합니다."""
    file_ext = file_ext.lower()
    for category, extensions in FILE_CATEGORIES.items():
        if category == '기타':
            continue
        if file_ext in extensions:
            return category
    return '기타'

def should_exclude_file(file_path: Path) -> bool:
    """파일이 제외 대상인지 확인합니다."""
    # 파일명 체크
    if file_path.name in EXCLUDE_PATTERNS:
        return True
    
    # 경로에 제외 패턴이 포함되어 있는지 체크
    for part in file_path.parts:
        if part in EXCLUDE_PATTERNS:
            return True
    
    return False

# --- AI 모듈 연동 함수 ---

def initialize_ai_module(model_path: str = None, prompt_template: str = None) -> AIKeywordExtractor | None:
    """
    AI 모듈을 초기화합니다.
    초기화는 프로그램 실행 중 한 번만 수행됩니다.
    """
    global AI_EXTRACTOR
    if AI_EXTRACTOR is not None:
        return AI_EXTRACTOR

    # 기본값 사용
    model_path = model_path or DEFAULT_MODEL_PATH
    prompt_template = prompt_template or DEFAULT_PROMPT_TEMPLATE

    model_path_obj = Path(model_path)
    if not model_path_obj.is_dir():
        print(f"오류: 모델 디렉토리를 찾을 수 없습니다. 경로를 확인하세요: {model_path_obj}")
        print(f"힌트: 현재 설정된 모델 경로는 '{DEFAULT_MODEL_PATH}' 입니다.")
        print("      다른 경로를 사용하려면 --model-path 옵션을 사용하세요.")
        return None

    print("🧠 Gemma ONNX 모델을 로딩합니다. 잠시 기다려주세요...")
    try:
        AI_EXTRACTOR = AIKeywordExtractor(
            model_path=str(model_path_obj),
            prompt_template=prompt_template
        )
        print("✅ AI 모델 로딩 완료!")
        return AI_EXTRACTOR
    except Exception as e:
        print(f"오류: AI 모델 로딩에 실패했습니다. ({e})")
        print("ONNX Runtime과 transformers 라이브러리가 올바르게 설치되었는지 확인해주세요.")
        return None


def get_ai_keywords_from_file(file_path: Path) -> str:
    """
    주어진 파일의 내용을 AI 모듈에 전달하여 핵심 키워드를 반환받습니다.
    """
    if AI_EXTRACTOR is None:
        return "AI_모듈_초기화_실패"

    print(f"🧠 AI: '{file_path.name}' 파일 내용 분석 중...")
    try:
        ext = file_path.suffix.lower()
        
        # 오디오 파일인 경우
        if ext in AI_SUPPORTED_AUDIO_EXTENSIONS:
            if not file_path.exists():
                return "파일없음"
            keywords = AI_EXTRACTOR.extract_keywords("", file_path.name, audio_path=str(file_path))
        # 이미지 파일인 경우
        elif ext in AI_SUPPORTED_IMAGE_EXTENSIONS:
            if not file_path.exists():
                return "파일없음"
            keywords = AI_EXTRACTOR.extract_keywords("", file_path.name, str(file_path))
        # 텍스트 파일인 경우
        elif ext in AI_SUPPORTED_EXTENSIONS:
            content = extract_text_from_file(file_path)
            if not content.strip():
                return "내용없음"
            keywords = AI_EXTRACTOR.extract_keywords(content, file_path.name)
        else:
            return "지원하지않는형식"
        
        # 파일명으로 부적절한 문자 제거
        safe_keywords = "".join(c for c in keywords if c.isalnum() or c in ['_', '-'])
        
        # 키워드가 너무 짧거나 의미없으면 파일명 기반으로 대체
        if not safe_keywords or safe_keywords in ["키워드추출실패", "내용없음", "지원하지않는형식", "오디오_분석불가_오디오모델없음"] or len(safe_keywords) < 3:
            # 파일명에서 의미있는 부분 추출
            name_parts = re.findall(r'[가-힣A-Za-z]+', file_path.stem)
            if name_parts:
                safe_keywords = '_'.join(name_parts[:3])  # 최대 3개 단어
            else:
                safe_keywords = "분류필요"
        
        return safe_keywords

    except Exception as e:
        print(f"오류: '{file_path.name}' 파일 분석 중 오류 발생: {e}")
        return "키워드추출실패"


def get_ai_filename_suggestion(file_path: Path) -> str:
    """
    AI를 사용하여 파일 내용 기반으로 새로운 파일명을 제안받습니다.
    """
    if AI_EXTRACTOR is None:
        return file_path.name

    print(f"🤖 AI: '{file_path.name}' 파일명 제안 중...")
    try:
        ext = file_path.suffix.lower()
        original_name = file_path.stem
        extension = file_path.suffix
        
        # 오디오 파일인 경우
        if ext in AI_SUPPORTED_AUDIO_EXTENSIONS:
            if not file_path.exists():
                return file_path.name
            suggested_name = AI_EXTRACTOR.suggest_filename("", original_name, extension, audio_path=str(file_path))
        # 이미지 파일인 경우
        elif ext in AI_SUPPORTED_IMAGE_EXTENSIONS:
            if not file_path.exists():
                return file_path.name
            suggested_name = AI_EXTRACTOR.suggest_filename("", original_name, extension, str(file_path))
        # 텍스트 파일인 경우
        elif ext in AI_SUPPORTED_EXTENSIONS:
            content = extract_text_from_file(file_path)
            if not content.strip():
                return file_path.name
            suggested_name = AI_EXTRACTOR.suggest_filename(content, original_name, extension)
        else:
            return file_path.name
        
        return suggested_name

    except Exception as e:
        print(f"오류: '{file_path.name}' 파일명 제안 중 오류 발생: {e}")
        return file_path.name


def get_ai_folder_classification(file_path: Path) -> str:
    """
    AI를 사용하여 파일 내용을 분석하고 적절한 폴더를 분류합니다.
    """
    if AI_EXTRACTOR is None:
        return "기타"

    print(f"📁 AI: '{file_path.name}' 폴더 분류 중...")
    try:
        ext = file_path.suffix.lower()
        
        # 오디오 파일인 경우
        if ext in AI_SUPPORTED_AUDIO_EXTENSIONS:
            if not file_path.exists():
                return "기타"
            folder_name = AI_EXTRACTOR.classify_folder("", file_path.name, audio_path=str(file_path))
        # 이미지 파일인 경우
        elif ext in AI_SUPPORTED_IMAGE_EXTENSIONS:
            if not file_path.exists():
                return "기타"
            folder_name = AI_EXTRACTOR.classify_folder("", file_path.name, str(file_path))
        # 텍스트 파일인 경우
        elif ext in AI_SUPPORTED_EXTENSIONS:
            content = extract_text_from_file(file_path)
            if not content.strip():
                return "기타"
            folder_name = AI_EXTRACTOR.classify_folder(content, file_path.name)
        else:
            return "기타"
        
        # 폴더명으로 부적절한 문자 제거
        safe_folder_name = "".join(c for c in folder_name if c.isalnum() or c in ['_', '-', ' '])
        return safe_folder_name if safe_folder_name else "기타"

    except Exception as e:
        print(f"오류: '{file_path.name}' 폴더 분류 중 오류 발생: {e}")
        return "기타"

# --- 핵심 로직 (Core Logic) ---

def organize_files(target_path: Path, model_path: str = None, dry_run: bool = False, 
                  recursive: bool = False, use_ai_keyword: bool = False, 
                  use_ai_filename: bool = False, use_ai_classify: bool = False, 
                  skip_confirm: bool = False, log_file: str = None):
    """파일 정리 작업을 수행합니다."""

    # AI 기능 사용 시, 모델 초기화
    if use_ai_keyword or use_ai_filename or use_ai_classify:
        if not initialize_ai_module(model_path):
            sys.exit(1) # 초기화 실패 시 프로그램 종료

    # 1. 정리할 파일 목록 스캔
    print(f"🔍 '{target_path.resolve()}' 경로를 스캔합니다...")
    files_to_organize = []
    file_iterator = target_path.rglob('*') if recursive else target_path.glob('*')

    for item in file_iterator:
        if item.is_file() and not should_exclude_file(item):
            files_to_organize.append(item)

    if not files_to_organize:
        print("정리할 파일이 없습니다.")
        return

    # 2. 변경 계획 수립
    change_plan = []
    for old_path in files_to_organize:
        new_name = old_path.name
        new_dir = old_path.parent

        original_name = old_path.stem
        ext = old_path.suffix.lower()
        date_created = datetime.datetime.fromtimestamp(old_path.stat().st_ctime).strftime('%Y-%m-%d')

        # 2-1. AI 기반 파일명 완전 변경
        if use_ai_filename and (ext in AI_SUPPORTED_EXTENSIONS or ext in AI_SUPPORTED_IMAGE_EXTENSIONS or ext in AI_SUPPORTED_AUDIO_EXTENSIONS):
            new_name = get_ai_filename_suggestion(old_path)
        
        # 2-2. AI 키워드 기반 이름 변경
        elif use_ai_keyword and (ext in AI_SUPPORTED_EXTENSIONS or ext in AI_SUPPORTED_IMAGE_EXTENSIONS or ext in AI_SUPPORTED_AUDIO_EXTENSIONS):
            keywords = get_ai_keywords_from_file(old_path)
            if keywords and keywords not in ["키워드추출실패", "내용없음", "지원하지않는형식", "오디오_분석불가_오디오모델없음"]:
                new_name = f"{original_name}_{keywords}{ext}"
        
        # 2-3. 일반 이름 변경 규칙 (텍스트 파일에만 날짜 추가)
        elif ext in ['.txt', '.md']:
            new_name = f"{date_created}_{original_name}{ext}"
        elif ext in ['.jpg', '.jpeg', '.png', '.gif']:
            new_name = f"IMG_{date_created}_{original_name}{ext}"
        elif ext in AI_SUPPORTED_AUDIO_EXTENSIONS:
            new_name = f"AUDIO_{date_created}_{original_name}{ext}"

        # 2-4. AI 기반 폴더 분류
        if use_ai_classify and (ext in AI_SUPPORTED_EXTENSIONS or ext in AI_SUPPORTED_IMAGE_EXTENSIONS or ext in AI_SUPPORTED_AUDIO_EXTENSIONS):
            ai_folder = get_ai_folder_classification(old_path)
            if ai_folder and ai_folder != "기타":
                new_dir = target_path / ai_folder
        
        # 2-5. 규칙 기반 폴더 이동 (AI 분류가 적용되지 않은 경우만)
        elif not use_ai_classify:
            category = get_file_category(ext)
            if category != '기타':
                new_dir = target_path / category
            else:
                new_dir = target_path / '기타'

        new_path = new_dir / new_name

        if old_path != new_path:
            change_plan.append({'from': old_path, 'to': new_path})

    # 3. 결과 출력 및 실행
    if not change_plan:
        print("✅ 모든 파일이 이미 정리되어 있습니다. 변경할 내용이 없습니다.")
        return

    print("\n✨ 다음과 같이 파일이 변경될 예정입니다:\n")
    for change in change_plan:
        try:
            from_rel = change['from'].relative_to(Path.cwd())
            to_rel = change['to'].relative_to(Path.cwd())
            print(f"  - [이동/변경] '{from_rel}' -> '{to_rel}'")
        except ValueError:
            print(f"  - [이동/변경] '{change['from']}' -> '{change['to']}'")

    if dry_run:
        print("\n💧 Dry-run 모드입니다. 실제 파일은 변경되지 않았습니다.")
        return

    # 4. 사용자 확인
    if not skip_confirm:
        confirm = input(f"\n총 {len(change_plan)}개의 파일을 변경하시겠습니까? (y/n): ")
        if confirm.lower() != 'y':
            print("작업을 취소했습니다.")
            return

    # 5. 실제 파일 변경 실행
    print("\n🚀 파일 정리를 시작합니다...")
    log_entries = []
    for change in change_plan:
        try:
            change['to'].parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(change['from']), str(change['to']))
            log_entry = f"[{datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] MOVED: '{change['from']}' -> '{change['to']}'"
            print(f"  - 성공: {change['from'].name} -> {change['to'].relative_to(target_path)}")
            log_entries.append(log_entry)
        except Exception as e:
            log_entry = f"[{datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] FAILED: '{change['from']}' -> '{change['to']}'. 이유: {e}"
            print(f"  - 실패: {change['from'].name}. 이유: {e}")
            log_entries.append(log_entry)

    # 6. 로그 기록
    if log_file:
        log_path = Path(log_file)
        with open(log_path, 'a', encoding='utf-8') as f:
            for entry in log_entries:
                f.write(entry + '\n')
        print(f"\n📝 작업 결과가 '{log_path.resolve()}' 파일에 기록되었습니다.")

    print("\n🎉 파일 정리가 완료되었습니다!")


def scan_files(target_path: Path, recursive: bool, by_ext: bool, by_date: bool):
    """폴더 현황을 분석하여 리포트를 출력합니다."""
    print(f"📊 '{target_path.resolve()}' 경로를 분석합니다...")
    
    files = []
    file_iterator = target_path.rglob('*') if recursive else target_path.glob('*')

    for item in file_iterator:
        if item.is_file():
            files.append(item)
    
    if not files:
        print("분석할 파일이 없습니다.")
        return

    print(f"\n총 {len(files)}개의 파일을 찾았습니다.")

    if by_ext:
        print("\n--- 확장자별 파일 현황 ---")
        ext_count = {}
        for file in files:
            ext = file.suffix.lower() or '.no_extension'
            ext_count[ext] = ext_count.get(ext, 0) + 1
        
        sorted_exts = sorted(ext_count.items(), key=lambda item: item[1], reverse=True)
        
        for ext, count in sorted_exts:
            print(f"  - {ext:<15}: {count}개")

    if by_date:
        print("\n--- 수정일 기준 파일 현황 (가장 오래된 파일 5개) ---")
        files.sort(key=lambda f: f.stat().st_mtime)
        for file in files[:5]:
            mtime = datetime.datetime.fromtimestamp(file.stat().st_mtime).strftime('%Y-%m-%d %H:%M')
            print(f"  - [{mtime}] {file.name}")

# --- CLI 인터페이스 (CLI Interface) ---

def main():
    parser = argparse.ArgumentParser(
        prog="aifo",
        description="🤖 AI 파일 정리 에이전트: AI와 규칙을 사용하여 파일을 자동으로 정리합니다.",
        epilog="자세한 사용법은 'aifo [COMMAND] --help'를 입력하세요."
    )
    subparsers = parser.add_subparsers(dest="command", required=True, help="수행할 작업")

    # 'organize' 명령어
    parser_organize = subparsers.add_parser("organize", help="규칙에 따라 파일을 정리합니다.")
    parser_organize.add_argument("target_path", type=Path, help="정리할 폴더 경로")
    parser_organize.add_argument("-m", "--model-path", type=str, default=DEFAULT_MODEL_PATH, 
                                help=f"AI 모델 디렉토리 경로 (기본값: {DEFAULT_MODEL_PATH})")
    parser_organize.add_argument("-d", "--dry-run", action="store_true", 
                                help="[안전 기능] 실제 파일 변경 없이 시뮬레이션 결과만 보여줍니다.")
    parser_organize.add_argument("-r", "--recursive", action="store_true", 
                                help="하위 폴더를 포함하여 전체를 정리합니다.")
    parser_organize.add_argument("-a", "--ai-keyword", action="store_true", 
                                help="[AI 기능] 문서 파일 내용 기반으로 키워드를 추출하여 파일명에 추가합니다.")
    parser_organize.add_argument("-f", "--ai-filename", action="store_true", 
                                help="[AI 기능] 문서 파일 내용을 분석하여 완전히 새로운 파일명을 제안합니다.")
    parser_organize.add_argument("-c", "--ai-classify", action="store_true", 
                                help="[AI 기능] 문서 파일 내용을 분석하여 적절한 폴더로 분류합니다.")
    parser_organize.add_argument("-A", "--ai-all", action="store_true", 
                                help="[AI 기능] 모든 AI 기능을 활성화합니다 (-a -f -c 와 동일).")
    parser_organize.add_argument("-y", "--yes", action="store_true", 
                                help="[주의] 파일 변경 전 확인 절차를 생략하고 즉시 실행합니다.")
    parser_organize.add_argument("-l", "--log", nargs='?', const=LOG_FILE_NAME, default=None, 
                                help=f"작업 로그를 파일로 기록합니다. 파일명을 지정하지 않으면 '{LOG_FILE_NAME}'에 저장됩니다.")

    # 'scan' 명령어
    parser_scan = subparsers.add_parser("scan", help="폴더 현황을 분석하여 리포트를 보여줍니다.")
    parser_scan.add_argument("target_path", type=Path, help="분석할 폴더 경로")
    parser_scan.add_argument("-r", "--recursive", action="store_true", help="하위 폴더를 포함하여 전체를 분석합니다.")
    parser_scan.add_argument("-e", "--by-ext", action="store_true", help="확장자별로 파일 개수를 보여줍니다.")
    parser_scan.add_argument("-t", "--by-date", action="store_true", help="수정일 기준으로 파일을 정렬하여 보여줍니다.")
    
    # 'config' 명령어
    parser_config = subparsers.add_parser("config", help="설정 정보를 확인합니다.")
    parser_config.add_argument("--show-categories", action="store_true", help="파일 카테고리 설정을 보여줍니다.")
    parser_config.add_argument("--show-ai-extensions", action="store_true", help="AI 지원 확장자를 보여줍니다.")
    parser_config.add_argument("--show-model-path", action="store_true", help="현재 모델 경로를 보여줍니다.")

    args = parser.parse_args()

    # AI 모든 기능 활성화 처리
    if hasattr(args, 'ai_all') and args.ai_all:
        args.ai_keyword = True
        args.ai_filename = True
        args.ai_classify = True

    # 명령어 실행
    if args.command == "organize":
        if not args.target_path.is_dir():
            print(f"오류: '{args.target_path}'는 유효한 디렉토리가 아닙니다.")
            sys.exit(1)
        
        organize_files(
            args.target_path, args.model_path, args.dry_run, args.recursive, 
            args.ai_keyword, args.ai_filename, args.ai_classify,
            args.yes, args.log
        )
    
    elif args.command == "scan":
        if not args.target_path.is_dir():
            print(f"오류: '{args.target_path}'는 유효한 디렉토리가 아닙니다.")
            sys.exit(1)
        # scan 에서는 by-ext나 by-date가 지정되지 않으면 둘 다 실행
        run_all = not args.by_ext and not args.by_date
        scan_files(args.target_path, args.recursive, args.by_ext or run_all, args.by_date or run_all)
        
    elif args.command == "config":
        if args.show_categories:
            print("📁 파일 카테고리 설정:")
            for category, extensions in FILE_CATEGORIES.items():
                if category == '기타':
                    print(f"  - {category}: 기타 모든 파일")
                else:
                    print(f"  - {category}: {', '.join(extensions)}")
        
        elif args.show_ai_extensions:
            print("🧠 AI 지원 확장자:")
            print(f"  - 텍스트: {', '.join(AI_SUPPORTED_EXTENSIONS)}")
            print(f"  - 이미지: {', '.join(AI_SUPPORTED_IMAGE_EXTENSIONS)}")
            print(f"  - 오디오: {', '.join(AI_SUPPORTED_AUDIO_EXTENSIONS)}")
        
        elif args.show_model_path:
            print(f"🤖 기본 모델 경로: {DEFAULT_MODEL_PATH}")
            model_path_obj = Path(DEFAULT_MODEL_PATH)
            if model_path_obj.exists():
                print("  ✅ 경로가 존재합니다.")
            else:
                print("  ❌ 경로가 존재하지 않습니다.")
        
        else:
            # 기본적으로 모든 설정 정보 표시
            print("📋 AIFO 설정 정보:")
            print(f"🤖 기본 모델 경로: {DEFAULT_MODEL_PATH}")
            print(f"🧠 AI 지원 확장자 - 텍스트: {', '.join(AI_SUPPORTED_EXTENSIONS)}")
            print(f"🧠 AI 지원 확장자 - 이미지: {', '.join(AI_SUPPORTED_IMAGE_EXTENSIONS)}")
            print(f"🧠 AI 지원 확장자 - 오디오: {', '.join(AI_SUPPORTED_AUDIO_EXTENSIONS)}")
            print(f"📁 파일 카테고리: {len(FILE_CATEGORIES)}개")
            print(f"🚫 제외 패턴: {', '.join(EXCLUDE_PATTERNS)}")


if __name__ == "__main__":
    main()